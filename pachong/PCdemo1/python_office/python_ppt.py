'''
https://blog.csdn.net/huuinn/article/details/78965725
'''
# 1.简单创建
# from pptx import Presentation
#
# # 创建幻灯片 ------
# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
#
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# # 设置标题和副标题
# title.text = "Hello, World!"
# subtitle.text = "pip install python-pptx"
# prs.save("test.pptx")
# 2.创建图标
# from pptx import Presentation
# from pptx.chart.data import ChartData
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.util import Inches
#
# # 创建幻灯片 ------
# prs = Presentation()
# slide = prs.slides.add_slide(prs.slide_layouts[5])
# # 幻灯片标题 ------
# title = slide.shapes.title
# title.text = "ppt表格"
# # 定义图表数据 ---------------------
# chart_data = ChartData()
# chart_data.categories = ['East', 'West', 'Midwest']
# chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
#
# # 将图表添加到幻灯片 --------------------
# x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
# slide.shapes.add_chart(
#     XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
# )
#
# prs.save('chart-01.pptx')

# 3.使用ppt模板来生成ppt

from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm  # Inches
from pptx.enum.chart import XL_LEGEND_POSITION

if __name__ == '__main__':
    # 创建幻灯片 ------
    prs = Presentation('template.pptx')
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = '报告'

    # 定义表格数据 ------
    name_objects = ["object1", "object2", "object3"]
    name_AIs = ["AI1", "AI2", "AI3"]
    val_AI1 = (19.2, 21.4, 16.7)
    val_AI2 = (22.3, 28.6, 15.2)
    val_AI3 = (20.4, 26.3, 14.2)
    val_AIs = [val_AI1, val_AI2, val_AI3]

    # 表格样式 --------------------
    rows = 4
    cols = 4
    top = Cm(12.5)
    left = Cm(3.5)  # Inches(2.0)
    width = Cm(24)  # Inches(6.0)
    height = Cm(6)  # Inches(0.8)

    # 添加表格到幻灯片 --------------------
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # 设置单元格宽度
    table.columns[0].width = Cm(6)  # Inches(2.0)
    table.columns[1].width = Cm(6)
    table.columns[2].width = Cm(6)
    table.columns[3].width = Cm(6)

    # 设置标题行
    table.cell(0, 1).text = name_objects[0]
    table.cell(0, 2).text = name_objects[1]
    table.cell(0, 3).text = name_objects[2]

    # 填充数据
    table.cell(1, 0).text = name_AIs[0]
    table.cell(1, 1).text = str(val_AI1[0])
    table.cell(1, 2).text = str(val_AI1[1])
    table.cell(1, 3).text = str(val_AI1[2])

    table.cell(2, 0).text = name_AIs[1]
    table.cell(2, 1).text = str(val_AI2[0])
    table.cell(2, 2).text = str(val_AI2[1])
    table.cell(2, 3).text = str(val_AI2[2])

    table.cell(3, 0).text = name_AIs[2]
    table.cell(3, 1).text = str(val_AI3[0])
    table.cell(3, 2).text = str(val_AI3[1])
    table.cell(3, 3).text = str(val_AI3[2])

    # 定义图表数据 ---------------------
    chart_data = ChartData()
    chart_data.categories = name_objects
    chart_data.add_series(name_AIs[0], val_AI1)
    chart_data.add_series(name_AIs[1], val_AI2)
    chart_data.add_series(name_AIs[2], val_AI3)

    # 添加图表到幻灯片 --------------------
    x, y, cx, cy = Cm(3.5), Cm(4.2), Cm(24), Cm(8)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False

    value_axis = chart.value_axis
    value_axis.maximum_scale = 100.0

    value_axis.has_title = True
    value_axis.axis_title.has_text_frame = True
    value_axis.axis_title.text_frame.text = "False positive"
    # value_axis.axis_title.text_frame.auto_size = 10

    prs.save('test_template.pptx')
